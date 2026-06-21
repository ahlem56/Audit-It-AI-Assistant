from __future__ import annotations

from io import BytesIO
from textwrap import shorten

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.models.export_models import ExportReportRequest

CONTROL_SHORT_LABELS = {
    "APD-01": "Révocation des accès",
    "APD-02": "Révocation des accès",
    "APD-03": "Comptes à droits étendus et génériques",
    "APD-04": "Recertification des droits d'accès",
    "APD-05": "Politique de sécurité des mots de passe",
    "PC-01": "Tests et validation avant mise en production",
    "PC-02": "Séparation des environnements",
    "CO-01": "Sauvegardes et plan de reprise",
    "CO-02": "Gestion des incidents de production",
    "CO-03": "Supervision des prestations externalisées",
    "CO-04": "Gestion des correctifs de sécurité",
    "CO-05": "Sauvegardes et plan de reprise",
    "CO-08": "Gestion des correctifs de sécurité",
}

PWC_RED = RGBColor(224, 48, 30)
PWC_ORANGE = RGBColor(235, 140, 0)
PWC_YELLOW = RGBColor(255, 182, 0)
PWC_DARK = RGBColor(45, 52, 65)
PWC_GREY = RGBColor(92, 103, 125)
PWC_LIGHT = RGBColor(247, 248, 250)
WHITE = RGBColor(255, 255, 255)


def _safe(value: object) -> str:
    return str(value or "").strip()


def _truncate(value: object, limit: int = 260) -> str:
    text = " ".join(_safe(value).split())
    return shorten(text, width=limit, placeholder="...")


def _priority_color(priority: str) -> RGBColor:
    normalized = priority.strip().lower()
    if normalized == "critical":
        return PWC_RED
    if normalized == "high":
        return PWC_ORANGE
    if normalized == "medium":
        return PWC_YELLOW
    return RGBColor(44, 151, 75)


def _control_label(finding) -> str:
    category = _safe(getattr(finding, "category", ""))
    if category:
        return category
    reference = _safe(getattr(finding, "reference", "")).upper()
    return CONTROL_SHORT_LABELS.get(reference, reference or "Contrôle non précisé")


def _control_application_label(finding) -> str:
    reference = _safe(getattr(finding, "reference", "")).upper()
    application = _safe(getattr(finding, "application", ""))
    label = _control_label(finding)
    first_line = f"{label} - {application}" if application else label
    return f"{first_line}\nRef. {reference}" if reference else first_line


def _add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int = 14,
    color: RGBColor = PWC_DARK,
    bold: bool = False,
    align: PP_ALIGN | None = None,
):
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    if align is not None:
        paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def _add_title(slide, title: str) -> None:
    _add_textbox(slide, Inches(0.55), Inches(0.35), Inches(11.7), Inches(0.6), title, size=24, bold=True)
    line = slide.shapes.add_shape(1, Inches(0.55), Inches(1.05), Inches(1.35), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = PWC_RED
    line.line.color.rgb = PWC_RED


def _add_footer(slide, label: str, slide_number: int) -> None:
    _add_textbox(slide, Inches(0.55), Inches(7.05), Inches(9.6), Inches(0.25), label, size=7, color=PWC_GREY)
    _add_textbox(slide, Inches(12.1), Inches(7.05), Inches(0.6), Inches(0.25), str(slide_number), size=7, color=PWC_GREY, align=PP_ALIGN.RIGHT)


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bullet_slide(prs: Presentation, title: str, bullets: list[str], footer: str) -> None:
    slide = _blank_slide(prs)
    _add_title(slide, title)
    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.45), Inches(11.4), Inches(5.25)).text_frame
    body.word_wrap = True
    body.clear()
    for index, bullet in enumerate(bullets[:9]):
        paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
        paragraph.text = _truncate(bullet, 210)
        paragraph.level = 0
        paragraph.font.name = "Arial"
        paragraph.font.size = Pt(15)
        paragraph.font.color.rgb = PWC_DARK
        paragraph.space_after = Pt(9)
    _add_footer(slide, footer, len(prs.slides))


def _section_slide(prs: Presentation, title: str, subtitle: str, footer: str) -> None:
    slide = _blank_slide(prs)
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(253, 239, 232)
    bg.line.fill.background()
    _add_textbox(slide, Inches(0.65), Inches(3.0), Inches(9.6), Inches(0.8), title, size=30, bold=True)
    _add_textbox(slide, Inches(0.68), Inches(3.85), Inches(9.6), Inches(0.45), subtitle, size=14, color=PWC_GREY)
    _add_footer(slide, footer, len(prs.slides))


def _summary_slide(prs: Presentation, data, footer: str) -> None:
    slide = _blank_slide(prs)
    _add_title(slide, "Executive summary")
    _add_textbox(slide, Inches(0.75), Inches(1.35), Inches(7.8), Inches(2.4), _truncate(data.executive_summary, 720), size=13)
    _add_textbox(slide, Inches(0.75), Inches(3.95), Inches(7.8), Inches(1.7), _truncate(data.general_synthesis, 470), size=12, color=PWC_GREY)

    counts = {item.priority: item.count for item in data.priority_summary}
    cards = [("Critical", counts.get("Critical", 0)), ("High", counts.get("High", 0)), ("Medium", counts.get("Medium", 0)), ("Low", counts.get("Low", 0))]
    for index, (label, value) in enumerate(cards):
        top = Inches(1.35 + index * 1.05)
        rect = slide.shapes.add_shape(1, Inches(9.15), top, Inches(2.75), Inches(0.78))
        rect.fill.solid()
        rect.fill.fore_color.rgb = _priority_color(label)
        rect.line.fill.background()
        _add_textbox(slide, Inches(9.35), top + Inches(0.13), Inches(1.35), Inches(0.28), label, size=10, color=WHITE, bold=True)
        _add_textbox(slide, Inches(10.95), top + Inches(0.08), Inches(0.6), Inches(0.34), str(value), size=18, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)
    _add_footer(slide, footer, len(prs.slides))


def _table_like_slide(prs: Presentation, title: str, rows: list[tuple[str, str, str]], footer: str) -> None:
    slide = _blank_slide(prs)
    _add_title(slide, title)
    top = Inches(1.45)
    for index, row in enumerate(rows[:7]):
        y = top + Inches(index * 0.72)
        fill = PWC_LIGHT if index % 2 == 0 else WHITE
        rect = slide.shapes.add_shape(1, Inches(0.65), y, Inches(12.0), Inches(0.58))
        rect.fill.solid()
        rect.fill.fore_color.rgb = fill
        rect.line.color.rgb = RGBColor(226, 232, 240)
        if title == "Consolidated action plan":
            _add_textbox(slide, Inches(0.85), y + Inches(0.08), Inches(3.3), Inches(0.42), _truncate(row[0], 58), size=8, bold=True)
            _add_textbox(slide, Inches(4.25), y + Inches(0.08), Inches(2.0), Inches(0.28), _truncate(row[1], 42), size=8)
            _add_textbox(slide, Inches(6.35), y + Inches(0.08), Inches(5.9), Inches(0.28), _truncate(row[2], 95), size=8, color=PWC_GREY)
        else:
            _add_textbox(slide, Inches(0.85), y + Inches(0.08), Inches(1.5), Inches(0.28), _truncate(row[0], 28), size=9, bold=True)
            _add_textbox(slide, Inches(2.45), y + Inches(0.08), Inches(4.0), Inches(0.28), _truncate(row[1], 70), size=9)
            _add_textbox(slide, Inches(6.6), y + Inches(0.08), Inches(5.7), Inches(0.28), _truncate(row[2], 95), size=9, color=PWC_GREY)
    _add_footer(slide, footer, len(prs.slides))


def _finding_slide(prs: Presentation, finding, footer: str) -> None:
    slide = _blank_slide(prs)
    priority = _safe(finding.priority) or "Priority"
    color = _priority_color(priority)
    _add_title(slide, _truncate(finding.title or finding.reference, 95))
    pill = slide.shapes.add_shape(1, Inches(0.75), Inches(1.25), Inches(2.1), Inches(0.38))
    pill.fill.solid()
    pill.fill.fore_color.rgb = color
    pill.line.fill.background()
    _add_textbox(slide, Inches(0.9), Inches(1.32), Inches(1.75), Inches(0.2), priority.upper(), size=8, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(3.05), Inches(1.27), Inches(8.8), Inches(0.45), _control_application_label(finding), size=9, color=PWC_GREY)

    sections = [
        ("Finding", finding.finding),
        ("Risk impact", finding.risk_impact),
        ("Recommendation", finding.recommendation),
    ]
    for index, (label, value) in enumerate(sections):
        left = Inches(0.75 + index * 4.1)
        card = slide.shapes.add_shape(1, left, Inches(2.0), Inches(3.75), Inches(3.75))
        card.fill.solid()
        card.fill.fore_color.rgb = PWC_LIGHT
        card.line.color.rgb = RGBColor(226, 232, 240)
        _add_textbox(slide, left + Inches(0.2), Inches(2.22), Inches(3.3), Inches(0.3), label, size=12, bold=True)
        _add_textbox(slide, left + Inches(0.2), Inches(2.72), Inches(3.3), Inches(2.6), _truncate(value, 430), size=10, color=PWC_GREY)
    _add_footer(slide, footer, len(prs.slides))


def build_report_pptx_docker(result: ExportReportRequest) -> BytesIO:
    data = result.structured_output
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    footer = f"{_safe(data.client_name)} | {_safe(data.report_period)}"

    cover = _blank_slide(prs)
    _add_textbox(cover, Inches(0.7), Inches(0.65), Inches(1.0), Inches(0.35), "pwc", size=22, color=PWC_DARK, bold=True)
    _add_textbox(cover, Inches(0.75), Inches(2.2), Inches(10.8), Inches(1.25), _safe(data.cover_title) or "Audit report", size=30, bold=True)
    _add_textbox(cover, Inches(0.78), Inches(3.55), Inches(9.5), Inches(0.45), _safe(data.cover_subtitle), size=15, color=PWC_GREY)
    _add_textbox(cover, Inches(0.78), Inches(6.65), Inches(7.0), Inches(0.25), _safe(data.confidentiality_notice), size=8, color=PWC_GREY)
    block_colors = [PWC_RED, PWC_ORANGE, PWC_YELLOW]
    for index, color in enumerate(block_colors):
        block = cover.shapes.add_shape(1, Inches(10.7 + index * 0.45), Inches(0.65 + index * 0.25), Inches(0.8), Inches(0.22))
        block.fill.solid()
        block.fill.fore_color.rgb = color
        block.line.fill.background()

    toc = data.table_of_contents or ["Context and approach", "Executive summary", "Detailed findings", "Action plan", "Appendices"]
    _bullet_slide(prs, "Table of contents", [f"{index}. {item}" for index, item in enumerate(toc, 1)], footer)
    _section_slide(prs, "Context and audit approach", "Scope, objectives and methodology", footer)
    _bullet_slide(prs, "Objectives", list(data.objectives or []) + [_safe(data.scope_summary)], footer)
    _bullet_slide(prs, "Applications in scope", list(data.applications or []) + list(data.covered_processes or []), footer)
    _summary_slide(prs, data, footer)
    _bullet_slide(prs, "Strategic priorities", list(data.strategic_priorities or []) + list(data.executive_highlights or []), footer)

    control_rows = [(item.reference, item.process, item.description) for item in data.covered_controls]
    for start in range(0, len(control_rows), 7):
        _table_like_slide(prs, "Covered controls", control_rows[start : start + 7], footer)

    _section_slide(prs, "Detailed findings", "Key observations and recommended actions", footer)
    prioritized = sorted(data.detailed_findings, key=lambda item: {"Critical": 0, "High": 1, "Medium": 2, "Low": 3}.get(item.priority, 4))
    for finding in prioritized[:12]:
        _finding_slide(prs, finding, footer)

    action_rows = [
        (
            _control_application_label(finding),
            finding.owner or finding.owners or "Owner to confirm",
            finding.immediate_action or finding.recommendation,
        )
        for finding in prioritized
    ]
    for start in range(0, len(action_rows), 7):
        _table_like_slide(prs, "Consolidated action plan", action_rows[start : start + 7], footer)

    _bullet_slide(prs, "Conclusion", [_safe(data.conclusion), _safe(data.maturity_assessment), _safe(data.priority_insight)], footer)

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output
